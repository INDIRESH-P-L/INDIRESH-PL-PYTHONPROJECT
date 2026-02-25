from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Create a new Presentation
prs = Presentation()

# Function to add a slide with title and content
def add_slide(prs, title_text, content_text):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    for point in content_text.split('\n'):
        if point.strip():
            p = tf.add_paragraph()
            p.text = point.strip('- ')
            p.level = 0
            p.font.size = Pt(24)

# Slide 1: Title Slide (Layout 0)
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Expense Tracker Application"
subtitle.text = "A modern, responsive, and robust personal finance management tool"

# Slide 2: Introduction
intro_text = """
- An interactive web application built with Python (Flask) and SQLite.
- Helps users track their daily income and expenses easily and efficiently.
- Provides a comprehensive dashboard summarizing financial health.
- Visualizes spending patterns through category-wise breakdown.
"""
add_slide(prs, "1. Introduction", intro_text)

# Slide 3: Problem Statement
problem_text = """
- Keeping track of daily spending is difficult using traditional notebooks or complex spreadsheets.
- Lack of immediate visibility into personal balance can lead to overspending.
- Users struggle to identify which categories (e.g., Food, Travel) consume most of their income.
- There is a need for a lightweight, accessible tool that prevents expenses from exceeding total income.
"""
add_slide(prs, "2. Problem Statement", problem_text)

# Slide 4: Innovation Ideas in the Project
innovation_text = """
- Real-time Income/Expense Guard: Automatically blocks users from adding an expense if it exceeds total income.
- Premium UI/UX: Custom-designed dashboard with dynamic color mapping, vibrant micro-animations, and glassmorphism.
- Zero-Configuration Setup: Completely self-contained SQLite database with auto-initialization and no complex server dependencies to install.
- Dynamic Visualizations: Auto-adjusting category bars (HTML/CSS-based) mapping out spending by importance and grouping.
"""
add_slide(prs, "3. Innovation Ideas in the Project", innovation_text)

# Slide 5: Software Used
software_text = """
- Backend Core: Python 3
- Web Framework: Flask
- Production Server: Waitress (WSGI server)
- Database: SQLite (Built-in standard library)
- Frontend Structure: HTML5
- Styling: Vanilla CSS (Custom modern design system, Google Fonts)
- Interactivity & API Calls: Vanilla JavaScript
"""
add_slide(prs, "4. Software Used", software_text)

# Slide 6: Number of Modules Used
modules_text = """
- Core Modules:
  - expencetrack.py (Main Application & Routing)
  - database.py (if strictly separated) / SQLite init
  - Templates Module (index.html for UI rendering)
- Python Standard Libraries:
  - os (file path management)
  - sqlite3 (embedded database engine)
  - datetime (timestamping transactions)
- Third-Party Packages:
  - flask (web server routing and templating)
  - waitress (production serving)
"""
add_slide(prs, "5. Number of Modules Used", modules_text)

# Save the presentation
output_path = os.path.join(os.path.dirname(__file__), "Expense_Tracker_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
